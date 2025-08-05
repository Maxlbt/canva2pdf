
import time
import os
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from PIL import Image
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches

# ---- Config ----
CANVA_URL = 'https://www.canva.com/design/DAGjlUtFq1Y/Dklp9PwjOtP5n5PJxNOs7w/view?utm_content=DAGjlUtFq1Y&utm_campaign=designshare&utm_medium=link2&utm_source=uniquelinks&utlId=he494a55727'  # 👉 Remplacer par ton lien Canva
NUM_SLIDES = 23# 👉 Remplacer
OUTPUT_FOLDER = 'screenshotsDeck'
PDF_FILENAME = 'Speeral_pitchdeck_investisseurs.pdf'# 👉 Remplacer
PPT_FILENAME = 'Speeral_pitchdeck_investisseurs.pptx'# 👉 Remplacer
SCREENSHOT_DELAY = 5

# ---- Setup Chrome in full screen ----
chrome_options = Options()
chrome_options.add_argument("--start-fullscreen")
chrome_options.add_argument("--disable-infobars")
chrome_options.add_argument("--disable-notifications")

driver = webdriver.Chrome(options=chrome_options)
driver.get(CANVA_URL)

time.sleep(8)  # Laisser le temps de charger la page Canva

# Forcer le mode plein écran de la présentation
body = driver.find_element(By.TAG_NAME, "body")
body.send_keys("f")  # Canva utilise souvent "f" pour le plein écran
time.sleep(5)

os.makedirs(OUTPUT_FOLDER, exist_ok=True)

# ---- Captures écran slide par slide ----
previous_shot = None
valid_screenshots = []
for i in range(NUM_SLIDES):
    time.sleep(SCREENSHOT_DELAY)
    filename = os.path.join(OUTPUT_FOLDER, f"slide_{i+1}.png")
    driver.save_screenshot(filename)
    print(f"Capture : {filename}")

    if previous_shot:
        current = Image.open(filename)
        if current.tobytes() == previous_shot.tobytes():
            print(f"⚠️ Doublon détecté à la slide {i+1}, suppression.")
            os.remove(filename)
            continue
    previous_shot = Image.open(filename)
    valid_screenshots.append(filename)

    body.send_keys(Keys.ARROW_RIGHT)

driver.quit()

# ---- Création du PDF ----
pdf = FPDF(orientation="L", unit="pt", format="A4")
for img_path in valid_screenshots:
    img = Image.open(img_path)
    img_width, img_height = img.size
    page_width, page_height = 842, 595
    aspect = img_width / img_height
    if page_width / page_height > aspect:
        new_height = page_height
        new_width = new_height * aspect
    else:
        new_width = page_width
        new_height = new_width / aspect
    x = (page_width - new_width) / 2
    y = (page_height - new_height) / 2
    pdf.add_page()
    pdf.image(img_path, x=x, y=y, w=new_width, h=new_height)
pdf.output(PDF_FILENAME)
print(f"✅ PDF généré : {PDF_FILENAME}")

# ---- Création PowerPoint ----
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
for img_path in valid_screenshots:
    slide = prs.slides.add_slide(blank_slide_layout)
    img = Image.open(img_path)
    img_width, img_height = img.size
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    aspect = img_width / img_height
    if slide_width / slide_height > aspect:
        height = slide_height
        width = height * aspect
    else:
        width = slide_width
        height = width / aspect
    left = (slide_width - width) / 2
    top = (slide_height - height) / 2
    slide.shapes.add_picture(img_path, left, top, width=width, height=height)
prs.save(PPT_FILENAME)
print(f"✅ PowerPoint généré : {PPT_FILENAME}")
